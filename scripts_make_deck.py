from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent
OUT = ROOT / 'ReliefKaki_pitch_deck_v2.pptx'
PHOTO = ROOT / 'src/assets/reliefkaki-hero.png'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

INK = RGBColor(25, 37, 42)
FOREST = RGBColor(37, 62, 57)
SAGE = RGBColor(220, 233, 218)
PAPER = RGBColor(247, 244, 238)
WHITE = RGBColor(255, 253, 248)
CLAY = RGBColor(220, 161, 126)
RUST = RGBColor(166, 83, 48)
MUTED = RGBColor(99, 112, 108)
LINE = RGBColor(210, 205, 194)
MONO = 'Aptos Mono'
SANS = 'Aptos Display'
SERIF = 'Georgia'


def rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font=SANS, align=None, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = value
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.space_after = Pt(0)
    if align is not None: p.alignment = align
    return box


def tag(slide, value, x=.62, y=.45, color=RUST):
    return text(slide, value.upper(), x, y, 6.7, .25, 9, color, True, MONO)


def divider(slide, y, x=.62, w=12.08, color=LINE):
    rect(slide, x, y, w, .012, color)


def number(slide, n, x, y, color=RUST):
    return text(slide, n, x, y, .5, .3, 10, color, True, MONO)


def quote(slide, value, x, y, w, h, color=INK):
    return text(slide, value, x, y, w, h, 25, color, False, SERIF)


def footer(slide, number_text):
    divider(slide, 7.02)
    text(slide, 'RELIEFKAKI', .62, 7.15, 2, .16, 7.5, MUTED, True, MONO)
    text(slide, number_text, 12.1, 7.15, .6, .16, 7.5, MUTED, True, MONO, PP_ALIGN.RIGHT)

# 01 / cover
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, FOREST)
slide.shapes.add_picture(str(PHOTO), Inches(7.3), Inches(0), width=Inches(6.033), height=Inches(7.5))
rect(slide, 7.3, 0, 2.8, 7.5, FOREST).fill.transparency = 32
text(slide, 'SPARKX⁺CHANGE · ALEXANDRA HOSPITAL CAREGIVER RESPITE', .72, .62, 5.9, .25, 9, CLAY, True, MONO)
text(slide, 'ReliefKaki\nBridge', .7, 1.34, 6.1, 1.75, 48, WHITE, True, SANS)
quote(slide, 'Help, on your own terms.', .75, 3.42, 5.3, .52, WHITE)
text(slide, 'A quiet, practical way for caregivers to offload one thing today — matched with trained youth volunteers, never clinical care.', .75, 4.42, 5.3, .68, 16, RGBColor(217, 228, 220))
rect(slide, .75, 5.62, 3.35, .48, CLAY, True)
text(slide, 'ONE SMALL ASK  ·  VISIBLE RELIEF', .95, 5.76, 3.0, .15, 9, FOREST, True, MONO)
text(slide, '01', .75, 6.75, .4, .2, 9, CLAY, True, MONO)

# 02 / tension
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, PAPER)
tag(slide, 'The tension')
text(slide, 'Caregivers know help exists.\nThey still do not use it.', .62, 1.05, 8.9, 1.32, 36, INK, True)
quote(slide, 'The gap is not information.\nIt is activation friction.', .63, 2.78, 4.4, .78, RUST)
text(slide, 'Shame · effort · cost · uncertainty · no clear next action', .65, 3.79, 5.4, .28, 12, MUTED)
rect(slide, 7.25, 1.12, 5.42, 4.9, WHITE)
text(slide, 'THE EVIDENCE', 7.65, 1.55, 2, .2, 9, RUST, True, MONO)
text(slide, '50.09%', 7.65, 2.02, 2.5, .63, 38, FOREST, True)
text(slide, 'of caregivers knew respite resources', 7.67, 2.72, 3.6, .3, 12, MUTED)
divider(slide, 3.33, 7.65, 4.5)
text(slide, '82.83%', 7.65, 3.66, 2.8, .63, 38, FOREST, True)
text(slide, 'of those aware had never used respite', 7.67, 4.35, 3.6, .3, 12, MUTED)
text(slide, 'SMU ROSA research · caregiver awareness is not activation', 7.65, 5.37, 4.25, .25, 8.5, RUST, False, MONO)
footer(slide, '02')

# 03 / design insight
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, SAGE)
tag(slide, 'The insight')
text(slide, 'Make help feel like\na task — not a confession.', .62, 1.05, 8.0, 1.35, 38, INK, True)
quote(slide, '“I need one thing done”\nis easier to say than\n“I need emotional support.”', .66, 3.05, 4.75, 1.35, INK)
for i, (head, body) in enumerate([
    ('PRACTICAL', 'Use solution-oriented language.'),
    ('PRIVATE', 'Let people choose silence and timing.'),
    ('PREDICTABLE', 'Show exactly what happens next.'),
]):
    x = 6.15 + i*2.16
    number(slide, f'0{i+1}', x, 2.97)
    text(slide, head, x, 3.36, 1.75, .22, 10, FOREST, True, MONO)
    text(slide, body, x, 3.75, 1.7, .8, 12, MUTED)
footer(slide, '03')

# 04 / solution
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, PAPER)
tag(slide, 'The response')
text(slide, 'ReliefKaki', .62, 1.03, 6.6, .58, 38, INK, True)
quote(slide, 'A silent-help task layer between “formal services exist” and “I need help tonight.”', .65, 1.78, 8.45, .65, RUST)
steps = [
    ('01', 'POST', 'One bounded non-clinical task.'),
    ('02', 'MATCH', 'A trained volunteer fits skill + time.'),
    ('03', 'COMPLETE', 'A receipt makes relief visible.'),
]
for i, (n, head, body) in enumerate(steps):
    x = .64 + i*4.14
    rect(slide, x, 3.15, 3.75, 2.1, WHITE)
    number(slide, n, x+.27, 3.47)
    text(slide, head, x+.27, 3.91, 2.6, .3, 18, FOREST, True)
    text(slide, body, x+.27, 4.46, 3.05, .45, 12, MUTED)
text(slide, 'CAREGIVER', .65, 5.84, 3.7, .2, 9, RUST, True, MONO)
text(slide, 'VOLUNTEER', 4.79, 5.84, 3.7, .2, 9, RUST, True, MONO)
text(slide, 'AH / C3U', 8.93, 5.84, 3.7, .2, 9, RUST, True, MONO)
footer(slide, '04')

# 05 / product moment
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, FOREST)
tag(slide, 'The product moment', .65, .5, CLAY)
text(slide, 'Silence is a setting.', .63, 1.02, 7.2, .62, 38, WHITE, True)
quote(slide, 'No call needed.\nJust message me.', .66, 2.05, 4.2, .86, WHITE)
text(slide, 'Silent Task lets a caregiver request practical help without a conversation, explanation, or pressure to perform gratitude.', .65, 3.48, 4.6, .78, 14, RGBColor(210, 222, 214))
rect(slide, 7.1, 1.08, 5.55, 4.95, WHITE)
text(slide, 'CAREGIVER REQUEST', 7.53, 1.48, 2.1, .2, 9, RUST, True, MONO)
text(slide, 'Pick up a simple dinner\nand leave it at the ward counter', 7.53, 1.97, 4.3, .66, 21, INK, True)
divider(slide, 2.95, 7.53, 4.2)
text(slide, 'ERRANDS', 7.53, 3.25, 1.3, .18, 9, MUTED, True, MONO)
text(slide, '50 points', 10.7, 3.25, 1.05, .18, 10, RUST, True, MONO)
rect(slide, 7.53, 3.84, 4.25, .82, SAGE)
text(slide, '●   SILENT TASK', 7.8, 4.05, 2.3, .18, 10, FOREST, True, MONO)
text(slide, 'NO CALL OR CONVERSATION EXPECTED', 7.8, 4.31, 3.2, .14, 7.8, MUTED, False, MONO)
rect(slide, 7.53, 5.06, 4.25, .48, FOREST, True)
text(slide, 'PLACE MY REQUEST   ↗', 7.53, 5.19, 4.25, .15, 9, WHITE, True, MONO, PP_ALIGN.CENTER)
footer(slide, '05')

# 06 / safety
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, SAGE)
tag(slide, 'Safety is the feature')
text(slide, 'Warmth needs\nclear edges.', .62, 1.05, 4.3, 1.16, 40, INK, True)
for i, (symbol, head, body, col) in enumerate([
    ('↳', 'WE DO', 'Errands, reminders, meals, wayfinding, forms, companionship and home organisation.', FOREST),
    ('×', 'WE DO NOT', 'Medication, wound care, lifting, diagnosis, personal care, medical interpretation or crisis support.', RUST),
    ('!', 'WE ESCALATE', 'Every concern has a named AH-linked owner, a clinical route and an accountable receipt.', FOREST),
]):
    x = 5.55 + i*2.38
    rect(slide, x, 2.07, 1.86, 2.95, PAPER)
    text(slide, symbol, x+.25, 2.4, .35, .3, 23, col, False, SERIF)
    text(slide, head, x+.25, 3.02, 1.35, .2, 10, INK, True, MONO)
    text(slide, body, x+.25, 3.47, 1.35, 1.15, 10.5, MUTED)
footer(slide, '06')

# 07 / supply operating system
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, PAPER)
tag(slide, 'Volunteer operating system')
text(slide, 'The supply side is\nnot assumed.', .62, 1.05, 6.5, 1.14, 38, INK, True)
text(slide, 'Leading volunteer-management patterns: clear opportunity scopes, screening, skill matching, shift / availability fit, verified hours and feedback loops.', .65, 2.63, 5.35, .72, 14, MUTED)
flow = [('1', 'RECRUIT', 'role fit'), ('2', 'BRIEF', 'safeguard'), ('3', 'TAG', 'skills'), ('4', 'MATCH', 'time + task'), ('5', 'RECOGNISE', 'VIA receipt')]
for i, (n, head, body) in enumerate(flow):
    x=.64+i*2.45
    number(slide, n, x, 4.45)
    text(slide, head, x, 4.84, 2.0, .2, 12, FOREST, True, MONO)
    text(slide, body, x, 5.18, 1.7, .22, 11, MUTED)
    if i < 4: text(slide, '→', x+1.82, 4.84, .24, .2, 13, RUST, False, MONO)
footer(slide, '07')

# 08 / pilot
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, CLAY)
tag(slide, 'Pilot: narrow on purpose', .62, .5, FOREST)
text(slide, 'Start small.\nLearn quickly.\nEarn trust.', .62, 1.04, 5.15, 1.7, 41, INK, True)
quote(slide, 'One discharge route.\nOne named operating lead.\nOne weekly rhythm.', .66, 3.34, 4.4, 1.15, INK)
for i, (n, head, body) in enumerate([
    ('0–2', 'WEEKS', '20 youth volunteers · safe task list · escalation briefing'),
    ('3–4', 'WEEKS', 'one route · daily moderation · completion review'),
    ('2–6', 'MONTHS', 'uptake · confidence · repeat use · incidents'),
]):
    x = 6.22
    y = 1.23 + i*1.72
    text(slide, n, x, y, .82, .35, 19, FOREST, True, MONO)
    text(slide, head, x+1.0, y+.04, 1.0, .2, 10, INK, True, MONO)
    text(slide, body, x+2.1, y+.02, 3.6, .4, 12, INK)
    divider(slide, y+.66, x, 5.5, RGBColor(190,132,101))
footer(slide, '08')

# 09 / standards
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, PAPER)
tag(slide, 'What leading systems teach us')
text(slide, 'Trust before choice.\nClarity before complexity.', .62, 1.05, 7.7, 1.12, 38, INK, True)
items = [
    ('CARE MARKETPLACES', 'Pre-filled task templates, visible trust signals and progressive disclosure reduce first-post friction.'),
    ('VOLUNTEER PLATFORMS', 'Skill + time matching, screening, verified service hours and coordinator dashboards make supply dependable.'),
    ('CAREGIVER SERVICES', 'Triage, clear boundaries, baseline / closure measures and referral back to professionals protect care quality.'),
]
for i, (head, body) in enumerate(items):
    y = 3.12+i*1.06
    number(slide, f'0{i+1}', .65, y)
    text(slide, head, 1.2, y, 2.65, .2, 10, RUST, True, MONO)
    text(slide, body, 4.1, y-.02, 7.95, .45, 12.5, MUTED)
footer(slide, '09')

# 10 / measurement
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, FOREST)
tag(slide, 'The scorecard', .62, .5, CLAY)
text(slide, 'Relief should be\nvisible.', .62, 1.05, 5.3, 1.1, 40, WHITE, True)
quote(slide, 'Measure the help\nthat actually lands.', .66, 2.76, 4.55, .75, WHITE)
metrics=[('TASKS', 'accepted + completed'), ('SILENT', 'low-contact uptake'), ('PEOPLE', 'repeat caregiver use'), ('SAFETY', 'clinical advice incidents = 0')]
for i,(head,body) in enumerate(metrics):
    x=6.15+(i%2)*3.1; y=1.5+(i//2)*2.14
    rect(slide,x,y,2.62,1.62,RGBColor(47, 77, 70))
    text(slide,head,x+.26,y+.32,2.0,.2,10,CLAY,True,MONO)
    text(slide,body,x+.26,y+.8,2.0,.35,13,WHITE,True)
footer(slide, '10')

# 11 / demo
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, PAPER)
tag(slide, 'What we built')
text(slide, 'A working pilot\nexperience — live.', .62, 1.05, 6.55, 1.14, 38, INK, True)
rect(slide, .64, 3.1, 5.0, 1.42, SAGE)
text(slide, 'POST A SILENT TASK', .95, 3.47, 2.7, .2, 10, RUST, True, MONO)
text(slide, 'Then claim, complete and record the help.', .95, 3.85, 3.95, .25, 14, FOREST, True)
text(slide, 'LIVE DEMO', 7.08, 1.64, 1.3, .18, 10, RUST, True, MONO)
text(slide, 'abelcjh.github.io/\nreliefkaki', 7.08, 2.1, 5.2, .7, 26, FOREST, True)
divider(slide, 3.2, 7.08, 4.7)
text(slide, 'BUILT INTO THE PROTOTYPE', 7.08, 3.56, 2.55, .16, 9, RUST, True, MONO)
text(slide, 'Silent Task · moderated task board · safe scope · volunteer recognition · pilot scorecard', 7.08, 3.98, 4.72, .75, 14, MUTED)
footer(slide, '11')

# 12 / close
slide = prs.slides.add_slide(prs.slide_layouts[6]); rect(slide, 0, 0, 13.333, 7.5, CLAY)
tag(slide, 'The ask', .62, .55, FOREST)
text(slide, 'Make respite as easy\nas posting one task.', .62, 1.08, 8.3, 1.25, 42, INK, True)
quote(slide, 'Silent help.\nVisible relief.', .66, 3.08, 4.5, .75, INK)
text(slide, 'To pilot well, we need:', 7.22, 3.0, 2.5, .2, 11, INK, True, MONO)
for i, item in enumerate(['an AH pilot owner', 'a safe task list + escalation SOP', 'a youth onboarding partner', 'a caregiver feedback loop']):
    y=3.5+i*.48
    text(slide, f'0{i+1}', 7.22, y, .42, .16, 9, RUST, True, MONO)
    text(slide, item, 7.78, y-.03, 3.4, .22, 13, INK, True)
text(slide, 'RELIEFKAKI', .65, 6.8, 2.5, .18, 9, FOREST, True, MONO)
text(slide, '12', 12.1, 6.8, .6, .18, 9, FOREST, True, MONO, PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT)
